from fastapi import FastAPI, UploadFile, File, Form
from transformers import pipeline
import pdfplumber
from pymongo import MongoClient
import io
from bson import ObjectId
import docx
from pptx import Presentation
import os
from dotenv import load_dotenv
from fastapi import HTTPException


app = FastAPI()
load_dotenv()
mongo_uri = os.getenv("MONGO_URL")

classifier = pipeline("zero-shot-classification", model="typeform/distilbert-base-uncased-mnli")

client = MongoClient(mongo_uri)

db = client["test"]
user_collection = db["users"]


def extract_text_from_pdf(file_bytes):
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
    return text.strip()


def extract_text_from_docx(file_bytes):
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text(file: UploadFile, file_bytes):
    ext = os.path.splitext(file.filename)[-1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_bytes)
    else:
        return None 


def get_user_labels(user_id: str):
    user = user_collection.find_one({"_id": ObjectId(user_id)})
    if not user or "className" not in user or not user["className"]:
        return []
    return user["className"]

@app.post("/classify")
async def classify_pdf(user_id: str = Form(...), file: UploadFile = File(...)):
    file_bytes = await file.read()
    text = extract_text(file, file_bytes)

    if not text:
        raise HTTPException(status_code=422, detail="Cannot extract text from the file.")

    labels = get_user_labels(user_id)
    if not labels:
        return HTTPException(status_code=404, detail="No classification labels found for user.")

    result = classifier(text, labels)
    best_label = result["labels"][0]
    confidence = result["scores"][0]

    if confidence <= 0.5:
        return HTTPException(status_code=400, detail="Low confidence. Cannot determine class.")

    return {
        "success": True,
        "predicted_class": best_label,
        "confidence": confidence
    }
